import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd

class DocumentParser:
    def parse_pdf(self, file_path: str) -> str:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text() + "\n"
        return text

    def parse_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        
        # Также извлечем текст из таблиц Word
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text.append(" | ".join(row_text))
                    
        return "\n".join(text)

    def parse_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide_idx, slide in enumerate(prs.slides):
            text.append(f"--- Слайд {slide_idx + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())
        return "\n".join(text)

    def parse_xlsx_to_facts(self, file_path: str) -> list[str]:
        """
        Преобразование строк таблиц Excel в текстовые утверждения (факты) для векторизации
        """
        facts = []
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            df = df.dropna(how='all')
            # Заполняем пропуски предыдущими значениями для заголовков (forward fill)
            df = df.ffill()
            for index, row in df.iterrows():
                row_dict = row.dropna().to_dict()
                if row_dict:
                    fact_str = f"Лист '{sheet_name}', строка {index + 1}: " + ", ".join([f"{k} = {v}" for k, v in row_dict.items()])
                    facts.append(fact_str)
        return facts

    def parse_file(self, file_path: str) -> tuple[str, list[str]]:
        """
        Парсинг файла по его расширению.
        Возвращает кортеж: (общий текст, список табличных фактов)
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self.parse_pdf(file_path), []
        elif ext == ".docx":
            return self.parse_docx(file_path), []
        elif ext == ".pptx":
            return self.parse_pptx(file_path), []
        elif ext in [".xlsx", ".xls"]:
            # Для таблиц текст пустой, возвращаем только список фактов
            return "", self.parse_xlsx_to_facts(file_path)
        else:
            raise ValueError(f"Неподдерживаемый формат файла: {ext}")

    def chunk_text(self, text: str, chunk_size: int = 600, overlap: int = 100) -> list[str]:
        words = text.split()
        chunks = []
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk:
                chunks.append(chunk)
        return chunks

parser = DocumentParser()

